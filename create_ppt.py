"""Generate the SmartScholar AI presentation deck."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color Palette ────────────────────────────────────────────────────────────
PRIMARY = RGBColor(0x66, 0x7E, 0xEA)
SECONDARY = RGBColor(0x76, 0x4B, 0xA2)
DARK = RGBColor(0x1E, 0x29, 0x3B)
GRAY = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)
ACCENT_BLUE = RGBColor(0x38, 0x82, 0xF6)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius:
        shape.adjustments[0] = radius
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide(slide, items, left, top, width, height, font_size=16, color=DARK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(10)
        p.level = 0
    return txBox


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide, WHITE)

add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PRIMARY)

add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1.2),
             "🎓 SmartScholar AI", font_size=48, color=PRIMARY, bold=True,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(0.8),
             "AI-Powered Research & Study Assistant", font_size=28, color=DARK,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(0.6),
             "Upload documents • Ask questions • Get intelligent, grounded answers",
             font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER)

add_shape_bg(slide, Inches(4.5), Inches(5.0), Inches(4.3), Inches(0.04), PRIMARY)

add_text_box(slide, Inches(1.5), Inches(5.4), Inches(10), Inches(0.5),
             "Vardhan Kumar K S  |  NeoStats AI Engineer Case Study",
             font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(5.9), Inches(10), Inches(0.5),
             "March 2026",
             font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2: Use Case Objective
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PRIMARY)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Use Case Objective", font_size=36, color=PRIMARY, bold=True)

add_shape_bg(slide, Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.04), LIGHT_BG)

add_text_box(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.8),
             "The Problem", font_size=24, color=DARK, bold=True)

add_bullet_slide(slide, [
    "• Students and researchers spend hours searching through multiple documents for answers",
    "• Traditional keyword search doesn't understand context or meaning",
    "• Information is often scattered across PDFs, notes, and the web",
    "• No single tool combines document Q&A with real-time web knowledge",
], Inches(0.8), Inches(2.4), Inches(11), Inches(2.0), font_size=17, color=GRAY)

add_text_box(slide, Inches(0.8), Inches(4.5), Inches(11), Inches(0.8),
             "Our Solution", font_size=24, color=DARK, bold=True)

add_bullet_slide(slide, [
    "• SmartScholar AI — an intelligent chatbot that understands your documents",
    "• Uses RAG (Retrieval-Augmented Generation) to ground answers in uploaded materials",
    "• Supplements with live web search for real-time, up-to-date information",
    "• Adapts response depth (Concise / Detailed) to match the user's needs",
], Inches(0.8), Inches(5.3), Inches(11), Inches(2.0), font_size=17, color=GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3: Approach
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PRIMARY)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Approach", font_size=36, color=PRIMARY, bold=True)

add_shape_bg(slide, Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.04), LIGHT_BG)

# Step boxes
steps = [
    ("1. Document Ingestion", "Upload PDF / DOCX / TXT files\n→ Extract text\n→ Split into overlapping chunks"),
    ("2. Vector Embedding", "Embed chunks using Google\nGemini Embedding API\n→ Store in FAISS index"),
    ("3. Semantic Retrieval", "User query embedded\n→ Cosine similarity search\n→ Top-K relevant chunks retrieved"),
    ("4. Augmented Generation", "Context + query sent to LLM\n→ Grounded, sourced response\n→ Streamed to user in real-time"),
]

for i, (title, desc) in enumerate(steps):
    x = Inches(0.6 + i * 3.15)
    box = add_shape_bg(slide, x, Inches(1.7), Inches(2.9), Inches(2.8), LIGHT_BG, radius=0.05)
    add_text_box(slide, x + Inches(0.2), Inches(1.85), Inches(2.5), Inches(0.5),
                 title, font_size=17, color=PRIMARY, bold=True)
    add_text_box(slide, x + Inches(0.2), Inches(2.4), Inches(2.5), Inches(2.0),
                 desc, font_size=14, color=GRAY)

add_text_box(slide, Inches(0.8), Inches(4.8), Inches(11), Inches(0.6),
             "Tech Stack", font_size=24, color=DARK, bold=True)

tech_items = [
    "• Frontend: Streamlit (Python web framework for data apps)",
    "• LLM Providers: Google Gemini 2.5 Flash / Groq Llama 3.3 / OpenAI GPT",
    "• Embeddings: Google Gemini Embedding API (768-dim vectors)",
    "• Vector Store: FAISS (Facebook AI Similarity Search) — in-memory, cosine similarity",
    "• Web Search: DuckDuckGo (free, no API key required)",
    "• Document Parsing: PyPDF2, python-docx",
]
add_bullet_slide(slide, tech_items, Inches(0.8), Inches(5.4), Inches(11), Inches(2.0),
                 font_size=15, color=GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4: Solution — Features
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PRIMARY)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Features Implemented", font_size=36, color=PRIMARY, bold=True)

add_shape_bg(slide, Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.04), LIGHT_BG)

features = [
    ("📄 RAG Integration", "Upload PDF, DOCX, TXT → chunked, embedded,\nindexed in FAISS → contextual answers\ngrounded in your documents"),
    ("🌐 Live Web Search", "Toggle DuckDuckGo search for real-time\nweb results when documents don't\nhave the answer"),
    ("💬 Response Modes", "Concise mode: quick 2-3 sentence answers\nDetailed mode: comprehensive, in-depth\nexplanations with examples"),
    ("🤖 Multi-LLM Support", "Switch between Google Gemini, Groq\n(Llama 3.3 70B), or OpenAI GPT\nwith a dropdown selector"),
    ("📋 Source Attribution", "Expandable panel showing which documents\nor web results informed each answer\nwith relevance scores"),
    ("⚡ Streaming Responses", "Real-time token-by-token output for\ninstant feedback — no waiting for\ncomplete generation"),
]

for i, (title, desc) in enumerate(features):
    col = i % 3
    row = i // 3
    x = Inches(0.6 + col * 4.2)
    y = Inches(1.6 + row * 2.7)
    box = add_shape_bg(slide, x, y, Inches(3.9), Inches(2.3), LIGHT_BG, radius=0.05)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(3.5), Inches(0.5),
                 title, font_size=18, color=DARK, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.7), Inches(3.5), Inches(1.5),
                 desc, font_size=14, color=GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5: Architecture
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PRIMARY)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Project Structure & Architecture", font_size=36, color=PRIMARY, bold=True)

add_shape_bg(slide, Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.04), LIGHT_BG)

structure = """project/
├── config/
│   └── config.py           ← API keys, model settings, constants
├── models/
│   ├── llm.py              ← Multi-provider LLM interface
│   └── embeddings.py       ← Embedding models for RAG
├── utils/
│   ├── document_processor.py ← PDF/DOCX/TXT extraction & chunking
│   ├── rag.py              ← FAISS vector store retrieval engine
│   └── web_search.py       ← DuckDuckGo live web search
├── app.py                  ← Main Streamlit application
└── requirements.txt"""

box = add_shape_bg(slide, Inches(0.6), Inches(1.6), Inches(6.0), Inches(5.2), LIGHT_BG, radius=0.03)
add_text_box(slide, Inches(0.9), Inches(1.8), Inches(5.5), Inches(5.0),
             structure, font_size=13, color=DARK, font_name="Courier New")

add_text_box(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(0.5),
             "Data Flow", font_size=22, color=DARK, bold=True)

flow_items = [
    "1️⃣  User uploads documents via sidebar",
    "2️⃣  Text extracted → split into 400-word chunks",
    "3️⃣  Chunks embedded → stored in FAISS index",
    "4️⃣  User asks a question",
    "5️⃣  Query embedded → top-4 similar chunks retrieved",
    "6️⃣  (Optional) DuckDuckGo web search runs",
    "7️⃣  Context + query → sent to chosen LLM",
    "8️⃣  Response streamed with source attribution",
]
add_bullet_slide(slide, flow_items, Inches(7.0), Inches(2.2), Inches(5.5), Inches(4.5),
                 font_size=16, color=GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6: Challenges
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PRIMARY)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Challenges & Solutions", font_size=36, color=PRIMARY, bold=True)

add_shape_bg(slide, Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.04), LIGHT_BG)

challenges = [
    ("Google SDK Deprecation",
     "The google-generativeai package was deprecated mid-development.",
     "Migrated to the new google-genai SDK with updated API patterns."),
    ("Model Availability & Rate Limits",
     "Gemini 2.0 Flash hit free tier quota limits (429 errors). Older models (1.5) were retired.",
     "Updated to Gemini 2.5 series and added multi-model selector for fallback options."),
    ("Package Naming Changes",
     "duckduckgo-search was renamed to ddgs, breaking Streamlit Cloud deployment.",
     "Updated imports and requirements to use the new package name."),
    ("API Key Security",
     "Pre-filled API keys in the UI could expose secrets to public app visitors.",
     "Keys from secrets are now hidden — only a confirmation badge is shown."),
]

for i, (title, challenge, solution) in enumerate(challenges):
    y = Inches(1.6 + i * 1.35)
    box = add_shape_bg(slide, Inches(0.6), y, Inches(12.0), Inches(1.2), LIGHT_BG, radius=0.03)
    add_text_box(slide, Inches(0.9), y + Inches(0.08), Inches(3.0), Inches(0.4),
                 title, font_size=16, color=PRIMARY, bold=True)
    add_text_box(slide, Inches(0.9), y + Inches(0.5), Inches(5.3), Inches(0.7),
                 "Challenge: " + challenge, font_size=13, color=GRAY)
    add_text_box(slide, Inches(6.5), y + Inches(0.5), Inches(5.8), Inches(0.7),
                 "Solution: " + solution, font_size=13, color=ACCENT_GREEN)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7: Demo & Deployment
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PRIMARY)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "Deployment & Links", font_size=36, color=PRIMARY, bold=True)

add_shape_bg(slide, Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.04), LIGHT_BG)

links = [
    ("🌐  Live App (Streamlit Cloud)", "https://smartscholar-ai-vardhan.streamlit.app"),
    ("📂  GitHub Repository", "https://github.com/Vardhankumarks/smartscholar-ai"),
]

for i, (label, url) in enumerate(links):
    y = Inches(1.8 + i * 1.4)
    box = add_shape_bg(slide, Inches(2.0), y, Inches(9.3), Inches(1.1), LIGHT_BG, radius=0.05)
    add_text_box(slide, Inches(2.4), y + Inches(0.1), Inches(8.5), Inches(0.45),
                 label, font_size=20, color=DARK, bold=True)
    add_text_box(slide, Inches(2.4), y + Inches(0.55), Inches(8.5), Inches(0.45),
                 url, font_size=16, color=PRIMARY)

add_text_box(slide, Inches(0.8), Inches(4.8), Inches(11), Inches(0.6),
             "Deployment Stack", font_size=22, color=DARK, bold=True)

deploy_items = [
    "• Hosted on Streamlit Cloud — auto-deploys from GitHub main branch",
    "• API keys stored securely in Streamlit Secrets (never exposed in UI)",
    "• No server management needed — fully serverless deployment",
    "• Each user gets an isolated session — documents are never shared between visitors",
]
add_bullet_slide(slide, deploy_items, Inches(0.8), Inches(5.4), Inches(11), Inches(2.0),
                 font_size=16, color=GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8: Thank You
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PRIMARY)

add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.0),
             "Thank You!", font_size=52, color=PRIMARY, bold=True,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(0.6),
             "SmartScholar AI — Where Technology Meets Knowledge",
             font_size=22, color=GRAY, alignment=PP_ALIGN.CENTER)

add_shape_bg(slide, Inches(4.5), Inches(4.2), Inches(4.3), Inches(0.04), PRIMARY)

add_text_box(slide, Inches(1.5), Inches(4.7), Inches(10), Inches(0.5),
             "Vardhan Kumar K S  •  Vardhanks342@gmail.com",
             font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(5.2), Inches(10), Inches(0.5),
             "github.com/Vardhankumarks/smartscholar-ai",
             font_size=14, color=PRIMARY, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════════════════
output = "/Users/dharshankumarks/Documents/vardhan_ai_engineer/SmartScholar_AI_Presentation.pptx"
prs.save(output)
print(f"Presentation saved to: {output}")
